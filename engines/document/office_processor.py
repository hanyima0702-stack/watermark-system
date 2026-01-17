"""
Office document processor for Word, Excel, and PowerPoint files.
Implements visible watermarks, invisible watermarks, and digital signatures.
"""
import os
import shutil
import tempfile
from pathlib import Path
from typing import Optional, Dict, Any, List
from datetime import datetime
import logging
import zipfile
import xml.etree.ElementTree as ET
from xml.dom import minidom

# Office document libraries
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import nsdecls, qn
from docx.oxml import parse_xml
import openpyxl
from openpyxl.drawing.image import Image as ExcelImage
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.oxml import parse_xml
from xml.sax.saxutils import escape
from pathlib import Path
from datetime import datetime
from typing import Dict, Any

from .base_processor import DocumentProcessor, WatermarkConfig, DocumentProcessingError
from .base_processor import WatermarkEmbeddingError, DigitalSignatureError

logger = logging.getLogger(__name__)


class OfficeProcessor(DocumentProcessor):
    """Processor for Microsoft Office documents (Word, Excel, PowerPoint)."""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.docx', '.xlsx', '.pptx']
        
    def add_visible_watermark(self, file_path: Path, watermark_config: WatermarkConfig,
                            user_id: str, timestamp: datetime) -> Path:
        """Add visible watermark to Office document."""
        if not self.is_supported_format(file_path):
            raise DocumentProcessingError(f"Unsupported format: {file_path.suffix}")
            
        output_path = self.generate_output_path(file_path)
        
        try:
            if file_path.suffix.lower() == '.docx':
                return self._add_word_watermark(file_path, output_path, watermark_config, user_id, timestamp)
            elif file_path.suffix.lower() == '.xlsx':
                return self._add_excel_watermark(file_path, output_path, watermark_config, user_id, timestamp)
            elif file_path.suffix.lower() == '.pptx':
                return self._add_powerpoint_watermark(file_path, output_path, watermark_config, user_id, timestamp)
        except Exception as e:
            logger.error(f"Failed to add visible watermark: {e}")
            raise WatermarkEmbeddingError(f"Failed to add visible watermark: {e}")
            
    def add_invisible_watermark(self, file_path: Path, watermark_data: str) -> Path:
        """Add invisible watermark using zero-width characters or document structure modification."""
        if not self.is_supported_format(file_path):
            raise DocumentProcessingError(f"Unsupported format: {file_path.suffix}")
            
        output_path = self.generate_output_path(file_path, "_invisible")
        
        try:
            if file_path.suffix.lower() == '.docx':
                return self._add_word_invisible_watermark(file_path, output_path, watermark_data)
            elif file_path.suffix.lower() == '.xlsx':
                return self._add_excel_invisible_watermark(file_path, output_path, watermark_data)
            elif file_path.suffix.lower() == '.pptx':
                return self._add_powerpoint_invisible_watermark(file_path, output_path, watermark_data)
        except Exception as e:
            logger.error(f"Failed to add invisible watermark: {e}")
            raise WatermarkEmbeddingError(f"Failed to add invisible watermark: {e}")
    
    def add_digital_signature(self, file_path: Path, certificate_path: Path, password: str) -> Path:
        """Add digital signature to Office document."""
        # Note: Full digital signature implementation requires COM components on Windows
        # This is a placeholder implementation that adds signature metadata
        output_path = self.generate_output_path(file_path, "_signed")
        
        try:
            shutil.copy2(file_path, output_path)
            
            # Add signature metadata to document properties
            if file_path.suffix.lower() == '.docx':
                self._add_word_signature_metadata(output_path, certificate_path)
            elif file_path.suffix.lower() == '.xlsx':
                self._add_excel_signature_metadata(output_path, certificate_path)
            elif file_path.suffix.lower() == '.pptx':
                self._add_powerpoint_signature_metadata(output_path, certificate_path)
                
            logger.info(f"Added signature metadata to {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Failed to add digital signature: {e}")
            raise DigitalSignatureError(f"Failed to add digital signature: {e}")
    
    def extract_invisible_watermark(self, file_path: Path) -> Optional[str]:
        """Extract invisible watermark from Office document."""
        try:
            if file_path.suffix.lower() == '.docx':
                return self._extract_word_invisible_watermark(file_path)
            elif file_path.suffix.lower() == '.xlsx':
                return self._extract_excel_invisible_watermark(file_path)
            elif file_path.suffix.lower() == '.pptx':
                return self._extract_powerpoint_invisible_watermark(file_path)
        except Exception as e:
            logger.error(f"Failed to extract invisible watermark: {e}")
            return None
    
    def _add_word_watermark(self, input_path: Path, output_path: Path, 
                          watermark_config: WatermarkConfig, user_id: str, timestamp: datetime) -> Path:
        """Add visible watermark to Word document."""
        doc = Document(str(input_path))
        
        # Get watermark configuration
        visible_config = watermark_config.visible_watermark
        if not visible_config.get('enabled', False):
            shutil.copy2(input_path, output_path)
            return output_path
            
        # Process each layer
        for layer in visible_config.get('layers', []):
            if layer.get('type') == 'text':
                self._add_word_text_watermark(doc, layer, user_id, timestamp)
                
        doc.save(str(output_path))
        return output_path

    def _add_word_text_watermark(self,doc: Document, layer_config: Dict[str, Any],
                                 user_id: str, timestamp: datetime):
        content = layer_config.get('content', '')
        content = content.replace('{user_id}', user_id)
        content = content.replace('{timestamp}', timestamp.strftime('%Y-%m-%d %H:%M:%S'))
        content = escape(content)  # 重要：XML 转义

        font = layer_config.get('font', {}) or {}
        position = layer_config.get('position', {}) or {}

        rotation = position.get('rotation', 0)
        margin_left = position.get('margin-left', 0)
        margin_top = position.get('margin-top', 0)
        x_pos = position.get('x', 'center')
        y_pos = position.get('y', 'center')

        family = font.get('family', 'Arial')
        width = font.get('width', 400)  # 水印通常较大
        height = font.get('height',100)
        color = font.get('color', '#808080')
        opacity = font.get('opacity', 0.5)
        # 把 opacity 转为字符串，VML 期望 0..1
        opacity_str = str(opacity)

        # 生成更规范的 VML，使用 <v:fill/> 控制颜色与透明度
        watermark_xml = f'''
        <w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"
             xmlns:v="urn:schemas-microsoft-com:vml"
             xmlns:o="urn:schemas-microsoft-com:office:office">
          <w:r>
            <w:pict>
              <v:shape id="PowerPlusWaterMarkObject" o:spid="_x0000_s2050" type="#_x0000_t136"
                       style="position:absolute;margin-left:{margin_left}pt;margin-top:{margin_top}pt;width:{width}pt;height:{height}pt;rotation:{rotation};z-index:-251658240;
                              mso-position-horizontal:{x_pos};mso-position-horizontal-relative:margin;
                              mso-position-vertical:{y_pos};mso-position-vertical-relative:margin"
                       o:allowincell="f" stroked="f">
                <v:fill color="{color}" opacity="{opacity_str}"/>
                <v:textpath on="t" style="font-family:'{family}'" string="{content}"/>
              </v:shape>
            </w:pict>
          </w:r>
        </w:p>
        '''

        # 将 watermark XML 注入每个节的 header（注意：解除与前一节的链接，以保证显示）
        for section in doc.sections:
            # 尝试断开与前一节的 header 链接（若需要）
            try:
                section.header.is_linked_to_previous = False
            except Exception:
                # 某些 python-docx 版本可能没有该属性，忽略
                pass

            header = section.header
            # 保证 header 至少有一个段落
            if not header.paragraphs:
                header.add_paragraph()
            # append XML 到 header 元素（使用 parse_xml）
            header._element.append(parse_xml(watermark_xml))
    
    def _add_excel_watermark(self, input_path: Path, output_path: Path,
                           watermark_config: WatermarkConfig, user_id: str, timestamp: datetime) -> Path:
        """Add visible watermark to Excel document."""
        workbook = openpyxl.load_workbook(str(input_path))
        
        visible_config = watermark_config.visible_watermark
        if not visible_config.get('enabled', False):
            workbook.save(str(output_path))
            return output_path
            
        # Add watermark to each worksheet
        for worksheet in workbook.worksheets:
            for layer in visible_config.get('layers', []):
                if layer.get('type') == 'text':
                    self._add_excel_text_watermark(worksheet, layer, user_id, timestamp)
                    
        workbook.save(str(output_path))
        return output_path
    
    def _add_excel_text_watermark(self, worksheet, layer_config: Dict[str, Any],
                                user_id: str, timestamp: datetime):
        """Add text watermark to Excel worksheet."""
        # Replace template variables
        content = layer_config.get('content', '')
        content = content.replace('{user_id}', user_id)
        content = content.replace('{timestamp}', timestamp.strftime('%Y-%m-%d %H:%M:%S'))
        
        # Add watermark text to header
        worksheet.oddHeader.center.text = content
        worksheet.evenHeader.center.text = content
    
    def _add_powerpoint_watermark(self, input_path: Path, output_path: Path,
                                watermark_config: WatermarkConfig, user_id: str, timestamp: datetime) -> Path:
        """Add visible watermark to PowerPoint document."""
        presentation = Presentation(str(input_path))
        
        visible_config = watermark_config.visible_watermark
        if not visible_config.get('enabled', False):
            presentation.save(str(output_path))
            return output_path
            
        # Add watermark to each slide
        for slide in presentation.slides:
            for layer in visible_config.get('layers', []):
                if layer.get('type') == 'text':
                    self._add_powerpoint_text_watermark(slide, layer, user_id, timestamp)
                    
        presentation.save(str(output_path))
        return output_path
    
    def _add_powerpoint_text_watermark(self, slide, layer_config: Dict[str, Any],
                                     user_id: str, timestamp: datetime):
        """Add text watermark to PowerPoint slide."""
        # Replace template variables
        content = layer_config.get('content', '')
        content = content.replace('{user_id}', user_id)
        content = content.replace('{timestamp}', timestamp.strftime('%Y-%m-%d %H:%M:%S'))
        
        # Add text box for watermark
        left = PptxInches(1)
        top = PptxInches(1)
        width = PptxInches(8)
        height = PptxInches(1)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = content
        
        # Apply font configuration
        font_config = layer_config.get('font', {})
        paragraph = text_frame.paragraphs[0]
        paragraph.font.name = font_config.get('family', 'Arial')
        paragraph.font.size = PptxPt(font_config.get('size', 12))
        
        # Set opacity
        opacity = font_config.get('opacity', 0.5)
        # Note: PowerPoint opacity requires more complex XML manipulation
    
    def _add_word_invisible_watermark(self, input_path: Path, output_path: Path, watermark_data: str) -> Path:
        """Add invisible watermark to Word document using zero-width characters."""
        doc = Document(str(input_path))
        
        # Convert watermark data to binary
        binary_data = ''.join(format(ord(char), '08b') for char in watermark_data)
        
        # Zero-width characters for binary encoding
        zero_width_space = '\u200B'  # 0
        zero_width_joiner = '\u200D'  # 1
        
        # Convert binary to zero-width characters
        invisible_text = ''
        for bit in binary_data:
            if bit == '0':
                invisible_text += zero_width_space
            else:
                invisible_text += zero_width_joiner
                
        # Insert invisible watermark at the beginning of the document
        if doc.paragraphs:
            first_paragraph = doc.paragraphs[0]
            if first_paragraph.runs:
                first_run = first_paragraph.runs[0]
                first_run.text = invisible_text + first_run.text
            else:
                first_paragraph.add_run(invisible_text)
        else:
            # Create new paragraph if document is empty
            paragraph = doc.add_paragraph()
            paragraph.add_run(invisible_text)
            
        doc.save(str(output_path))
        return output_path
    
    def _add_excel_invisible_watermark(self, input_path: Path, output_path: Path, watermark_data: str) -> Path:
        """Add invisible watermark to Excel document using custom properties."""
        workbook = openpyxl.load_workbook(str(input_path))
        
        # Add watermark as custom property
        workbook.properties.custom = workbook.properties.custom or {}
        
        # Encode watermark data
        import base64
        encoded_data = base64.b64encode(watermark_data.encode()).decode()
        
        # Store in custom property with innocuous name
        workbook.custom_doc_props = getattr(workbook, 'custom_doc_props', {})
        workbook.custom_doc_props['DocumentVersion'] = encoded_data
        
        workbook.save(str(output_path))
        return output_path
    
    def _add_powerpoint_invisible_watermark(self, input_path: Path, output_path: Path, watermark_data: str) -> Path:
        """Add invisible watermark to PowerPoint document using slide notes."""
        presentation = Presentation(str(input_path))
        
        # Encode watermark data
        import base64
        encoded_data = base64.b64encode(watermark_data.encode()).decode()
        
        # Add to first slide's notes (if exists)
        if presentation.slides:
            slide = presentation.slides[0]
            if hasattr(slide, 'notes_slide') and slide.notes_slide:
                notes_text_frame = slide.notes_slide.notes_text_frame
                # Hide watermark in notes as invisible characters
                zero_width_space = '\u200B'
                notes_text_frame.text = zero_width_space + encoded_data + zero_width_space + notes_text_frame.text
                
        presentation.save(str(output_path))
        return output_path
    
    def _extract_word_invisible_watermark(self, file_path: Path) -> Optional[str]:
        """Extract invisible watermark from Word document."""
        try:
            doc = Document(str(file_path))
            
            # Look for zero-width characters in first paragraph
            if doc.paragraphs:
                first_paragraph = doc.paragraphs[0]
                text = first_paragraph.text
                
                # Extract zero-width characters
                zero_width_space = '\u200B'
                zero_width_joiner = '\u200D'
                
                binary_data = ''
                for char in text:
                    if char == zero_width_space:
                        binary_data += '0'
                    elif char == zero_width_joiner:
                        binary_data += '1'
                        
                if binary_data and len(binary_data) % 8 == 0:
                    # Convert binary to text
                    watermark_data = ''
                    for i in range(0, len(binary_data), 8):
                        byte = binary_data[i:i+8]
                        watermark_data += chr(int(byte, 2))
                    return watermark_data
                    
        except Exception as e:
            logger.error(f"Failed to extract Word invisible watermark: {e}")
            
        return None
    
    def _extract_excel_invisible_watermark(self, file_path: Path) -> Optional[str]:
        """Extract invisible watermark from Excel document."""
        try:
            workbook = openpyxl.load_workbook(str(file_path))
            
            # Check custom properties
            if hasattr(workbook, 'custom_doc_props') and workbook.custom_doc_props:
                encoded_data = workbook.custom_doc_props.get('DocumentVersion')
                if encoded_data:
                    import base64
                    try:
                        watermark_data = base64.b64decode(encoded_data).decode()
                        return watermark_data
                    except Exception:
                        pass
                        
        except Exception as e:
            logger.error(f"Failed to extract Excel invisible watermark: {e}")
            
        return None
    
    def _extract_powerpoint_invisible_watermark(self, file_path: Path) -> Optional[str]:
        """Extract invisible watermark from PowerPoint document."""
        try:
            presentation = Presentation(str(file_path))
            
            # Check first slide's notes
            if presentation.slides:
                slide = presentation.slides[0]
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text
                    
                    # Extract encoded data between zero-width spaces
                    zero_width_space = '\u200B'
                    if zero_width_space in notes_text:
                        parts = notes_text.split(zero_width_space)
                        if len(parts) >= 3:
                            encoded_data = parts[1]
                            import base64
                            try:
                                watermark_data = base64.b64decode(encoded_data).decode()
                                return watermark_data
                            except Exception:
                                pass
                                
        except Exception as e:
            logger.error(f"Failed to extract PowerPoint invisible watermark: {e}")
            
        return None
    
    def _add_word_signature_metadata(self, file_path: Path, certificate_path: Path):
        """Add signature metadata to Word document."""
        # This is a simplified implementation
        # Full digital signature requires COM components
        doc = Document(str(file_path))
        
        # Add signature information to document properties
        core_props = doc.core_properties
        core_props.author = "Digital Signature System"
        core_props.comments = f"Digitally signed with certificate: {certificate_path.name}"
        core_props.last_modified_by = "Watermark System"
        
        doc.save(str(file_path))
    
    def _add_excel_signature_metadata(self, file_path: Path, certificate_path: Path):
        """Add signature metadata to Excel document."""
        workbook = openpyxl.load_workbook(str(file_path))
        
        # Add signature information to properties
        workbook.properties.creator = "Digital Signature System"
        workbook.properties.description = f"Digitally signed with certificate: {certificate_path.name}"
        
        workbook.save(str(file_path))
    
    def _add_powerpoint_signature_metadata(self, file_path: Path, certificate_path: Path):
        """Add signature metadata to PowerPoint document."""
        presentation = Presentation(str(file_path))
        
        # Add signature information to properties
        core_props = presentation.core_properties
        core_props.author = "Digital Signature System"
        core_props.comments = f"Digitally signed with certificate: {certificate_path.name}"
        
        presentation.save(str(file_path))